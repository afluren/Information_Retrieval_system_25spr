# FIXME
# 1.有些链接是引用了nankai.edu.cn的，但是并没有真正的链接，需要进一步处理


import requests
from bs4 import BeautifulSoup
import os
from urllib.parse import urljoin, unquote
import json
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import tempfile
import re
from multiprocessing import Process, Manager
from typing import List, Dict
import time
import pickle
import gzip
import base64
import chardet

# if res.status_code == 200:
#     html = res.text
#     soup = BeautifulSoup(html, 'html.parser')
#     a_tags = soup.find_all('a')
#     for a in a_tags:
#         title = a.get('title')
#         url = a.get('href')
#         if title:
#             print(url+": "+title)
#         else:
#             print(url+": "+a.text)
# else:
#     print(f"Failed to retrieve the page. Status code: {res.status_code}")


def extract_webpage(url: str, res: requests.Response, visited: Dict[str, bool], will_visit: Dict[str, bool], saved_num: int, lock, index):
    BASE = "nankai.edu.cn"
    # 爬取网页
    # res = requests.get(url, headers=headers)
    # res.encoding = res.encoding or 'utf-8'
    sample = res.content[:20000]
    res.encoding = chardet.detect(sample)['encoding'] or 'utf-8'
    if res.status_code == 200:
        website = {}  # 存储网页信息
        hrefs = set()  # 存储链接
        html = res.text  # 原始HTML
        soup = BeautifulSoup(html, 'html.parser')
        # 清洗富文本
        for tag in soup(['img', 'video', 'audio', 'source', 'iframe', 'embed', 'object', 'picture']):
            tag.decompose()
        for tag in soup.find_all(['iframe', 'embed', 'object']):
            src = tag.get('src') or tag.get('data') or ''
            if any(ext in src.lower() for ext in ['.pdf', '.doc', '.docx', '.ppt', '.pptx', '.xls', '.xlsx']):
                tag.decompose()

        # 清理网页内容
        pure_text = soup.get_text(separator='\n', strip=True)

        if '�' in pure_text or any(ord(c) < 32 and c not in '\n\r\t' for c in pure_text):
            raise ValueError(f"网页含有非法字符，直接跳过：{url}")
        a_tags = soup.find_all('a')
        title = soup.title.string if soup.title else ""
        links = []  # 清理后的anchors
        for a in a_tags:
            a_href = a.get('href')
            a_text = a.get('title')
            if a_text == None:
                a_text = a.text
            a_url = urljoin(url, a_href)
            if BASE in a_url and a_url.startswith("http"):
                if a_url not in hrefs:
                    hrefs.add(a_url)
                    links.append({"text": a_text, "href": a_url})

        # 压缩HTML
        html_zip = base64.b64encode(gzip.compress(res.content)).decode('ascii')

        website['text'] = pure_text
        website['title'] = title
        website['url'] = url
        website['anchors'] = links
        website['html_raw_compressed'] = html_zip
        # 保存网页信息
        with open(f'raw_html/nankai{index}.jsonl', 'a', encoding='utf-8') as f:
            json_line = json.dumps(website, ensure_ascii=False)
            f.write(json_line + '\n')
            print(f"Saved {url} to nankai{index}.jsonl")
        with lock:
            saved_num.value += 1
            # 递归爬取链接
            for next_url in hrefs:
                if next_url not in visited:
                    will_visit[next_url] = True

    else:
        print(
            f"Failed to retrieve the page. Status code: {res.status_code}, URL: {url}")


def extract_file_text(path, file_name):
    print(f"Extracting text from file: {path}")
    try:
        if file_name.lower().endswith('.pdf'):
            return extract_text(path)
        elif file_name.lower().endswith('.docx'):
            doc = Document(path)
            return '\n'.join(p.text for p in doc.paragraphs)
        elif file_name.lower().endswith('.pptx'):
            prs = Presentation(path)
            return '\n'.join(
                shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
            )
        elif file_name.lower().endswith('.txt'):
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif file_name.lower().endswith('.xlsx'):
            wb = load_workbook(path, read_only=True, data_only=True)
            text = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            text.append(str(cell.value))
            wb.close()
            return '\n'.join(text)
        else:
            return None
    except Exception as e:
        print(f"Failed to extract text from {path}: {e}")
        return None


def extract_file(url: str, r: requests.Response, visited: Dict[str, bool], saved_num: int, lock, index):
    print(f"Extracting file: {url}")
    # r = requests.get(url, timeout=10)  # 发请求下载文件
    if r.status_code == 200:
        # 获取文件名
        cd = r.headers.get("Content-Disposition", "")
        fname_match = re.search(r'filename="?([^"]+)"?', cd)
        if fname_match:
            file_name = unquote(fname_match.group(1))
        else:
            file_name = unquote(url.split("/")[-1])

        print(f"Downloading file: {file_name}")
        ext = os.path.splitext(file_name)[1]     # 获取后缀，如 '.pdf'
        fd, temp_path = tempfile.mkstemp(suffix=ext)  # 创建临时文件
        with os.fdopen(fd, 'wb') as f:
            f.write(r.content)             # 写入文件内容（二进制）
        file_text = extract_file_text(temp_path, file_name)  # 解析文件内容
        if os.path.exists(temp_path):
            os.remove(temp_path)           # 删除临时文件
        if file_text == None:
            return
        website = {
            'url': url,
            'type': "file",
            'filetype': os.path.splitext(url)[1][1:],
            'title': file_name,
            'text': file_text.strip(),
            'html_raw': None,
            'anchors': []
        }
        # 保存网页文件信息
        with open(f'raw_html/nankai{index}.jsonl', 'a', encoding='utf-8') as f:
            json_line = json.dumps(website, ensure_ascii=False)
            f.write(json_line + '\n')
            print(f"Saved {url} to nankai{index}.jsonl")
        with lock:
            saved_num.value += 1

    else:
        print(
            f"Failed to retrieve the file. Status code: {r.status_code}, URL: {url}")


def extract_links(url: str, visited, will_visit, saved_num, lock, index):
    # if url.startswith("http"):
    #     if url.lower().endswith(('.pdf', '.docx', '.pptx', '.xlsx', '.txt')):
    #         extract_file(url)
    #     elif url.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.svg', '.ico', '.bmp', '.tif')):
    #         print(f"Skipping image: {url}")
    #     else:
    #         extract_webpage(url)
    # else:
    #     print(f"Invalid URL: {url}")
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36"
    }
    retry = True
    for _ in range(2):
        if not retry:
            break
        with lock:
            visited[url] = True
        try:
            r = requests.get(url, headers=headers, timeout=(3, 3))
        except Exception as e:
            print(f"Failed to retrieve the page: {e}, URL: {url}")
            if url.startswith("https"):
                url = url.replace("https", "http", 1)
            else:
                url = url.replace("http", "https", 1)
            continue
        if r.status_code == 200:
            cd = r.headers.get("Content-Disposition", "").lower()
            ct = r.headers.get("Content-Type", "").lower()

            # 是否是“下载型”资源
            is_file = "attachment" in cd or "filename=" in cd or ct.startswith(
                "application/") or ct == "text/plain"
            if is_file:
                extract_file(url, r, visited, saved_num, lock, index)
            else:
                extract_webpage(url, r, visited, will_visit,
                                saved_num, lock, index)
            retry = False
        else:
            print(
                f"Failed to retrieve the page. Status code: {r.status_code}, URL: {url}")
            if url.startswith("https"):
                url = url.replace("https", "http", 1)
            else:
                url = url.replace("http", "https", 1)


def main(visited, will_visit: Dict[str, bool], saved_num, lock, index):
    URL_NUMBER = 300000
    while True:
        with lock:
            if saved_num.value >= URL_NUMBER:
                print(f"Reached the limit of {URL_NUMBER} pages. Stopping.")
                break
            if saved_num.value % 10000 == 0:
                with open("crawl_state.pkl", "wb") as f:
                    pickle.dump((dict(visited), dict(
                        will_visit), saved_num.value), f)
            if len(will_visit) == 0:
                break
            else:
                url = will_visit.popitem()[0]
        try:
            extract_links(url, visited, will_visit, saved_num, lock, index)
        except Exception as e:
            print(f"Failed to extract links from {url}: {e}")


if __name__ == '__main__':
    test_urls = ["https://www.nankai.edu.cn/", "https://wxy.nankai.edu.cn/", "https://history.nankai.edu.cn/", "https://phil.nankai.edu.cn/", "https://sfs.nankai.edu.cn/",
                 "https://law.nankai.edu.cn/", "https://zfxy.nankai.edu.cn/", "https://cz.nankai.edu.cn/", "https://hyxy.nankai.edu.cn/", "https://economics.nankai.edu.cn/", "https://cc.nankai.edu.cn/", "https://bs.nankai.edu.cn/", "https://tas.nankai.edu.cn/", "https://finance.nankai.edu.cn/", "https://math.nankai.edu.cn/", "https://physics.nankai.edu.cn/", "https://chem.nankai.edu.cn/", "https://shxy.nankai.edu.cn/", "https://jc.nankai.edu.cn/", "https://stat.nankai.edu.cn/", "https://cs.nankai.edu.cn/", "https://ai.nankai.edu.cn/", "https://cyber.nankai.edu.cn/", "https://mse.nankai.edu.cn/", "https://ceo.nankai.edu.cn/", "https://pharmacy.nankai.edu.cn/", "https://medical.nankai.edu.cn/", "https://env.nankai.edu.cn/", "https://sky.nankai.edu.cn/", "https://aiguo.nankai.edu.cn/", "https://xs.nankai.edu.cn/", "https://hq.nankai.edu.cn/"]
    # URL_NUMBER = 200
    # BASE = "nankai.edu.cn"
    start_time = time.time()
    with Manager() as manager:
        visited = manager.dict()
        will_visit = manager.dict()
        for test_url in test_urls:
            will_visit[test_url] = True
        saved_num = manager.Value('i', 0)
        lock = manager.Lock()
        num_workers = 32
        process = []
        reload_state = False
        while True:
            if reload_state:
                with open("crawl_state.pkl", "rb") as f:
                    visited2, will_visit2, saved_num2 = pickle.load(f)
                visited = manager.dict(visited2)
                will_visit = manager.dict(will_visit2)
                saved_num.value = saved_num2
            try:
                for i in range(num_workers):
                    p = Process(target=main, args=(
                        visited, will_visit, saved_num, lock, i))
                    p.start()
                    process.append(p)
                for p in process:
                    p.join()
                break
            except Exception as e:
                print(f"Failed to start worker: {e}")
                with open("crawl_state.pkl") as f:
                    pickle.dump((dict(visited), dict(
                        will_visit), saved_num.value), f)
                reload_state = True

    end_time = time.time()
    print(f"Time elapsed: {end_time - start_time} seconds")
